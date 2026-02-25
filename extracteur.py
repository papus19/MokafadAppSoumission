"""
Extracteur de texte universel — supporte PDF, Word, Excel, PowerPoint, images, texte brut.
Utilisé par analyse.py et generateur_offres.py.
"""
import io
from typing import List
import streamlit as st


# ── Types acceptés par les file_uploader ──────────────────────────────
TYPES_ACCEPTES = None  # None = tous les types dans Streamlit

# Label et help uniformes pour tous les uploaders multi-documents
LABEL_UPLOAD = "Glissez vos documents ici ou cliquez pour sélectionner"
HELP_UPLOAD  = (
    "Tous types acceptés : PDF, Word (.docx/.doc), Excel (.xlsx/.xls), "
    "PowerPoint (.pptx), images (.png/.jpg/.tiff), texte (.txt/.csv) — "
    "Plusieurs fichiers simultanément"
)


def _extraire_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    pages  = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extraire_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc   = Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # Tables
    for table in doc.tables:
        for row in table.rows:
            paras.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
    return "\n".join(paras)


def _extraire_xlsx(file_bytes: bytes) -> str:
    import openpyxl
    wb    = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lignes = []
    for sheet in wb.worksheets:
        lignes.append(f"=== Feuille : {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None]
            if vals:
                lignes.append(" | ".join(vals))
    return "\n".join(lignes)


def _extraire_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs   = Presentation(io.BytesIO(file_bytes))
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        texts.append(f"=== Diapositive {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def _extraire_image(file_bytes: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
        img  = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img, lang="fra+eng")
        return text.strip()
    except Exception as e:
        return f"[Image — extraction OCR échouée : {e}]"


def _extraire_txt(file_bytes: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def extraire_texte(uploaded_file) -> str:
    """
    Extrait le texte brut d'un fichier Streamlit UploadedFile.
    Supporte : pdf, docx, doc, xlsx, xls, pptx, png, jpg, jpeg, tiff, bmp, txt, csv.
    Retourne une chaîne vide avec un message d'erreur si le format n'est pas supporté.
    """
    nom  = uploaded_file.name.lower()
    data = uploaded_file.read()
    # Rembobiner pour usage ultérieur éventuel
    uploaded_file.seek(0)

    ext = nom.rsplit(".", 1)[-1] if "." in nom else ""

    try:
        if ext == "pdf":
            return _extraire_pdf(data)
        elif ext in ("docx", "doc"):
            return _extraire_docx(data)
        elif ext in ("xlsx", "xls"):
            return _extraire_xlsx(data)
        elif ext == "pptx":
            return _extraire_pptx(data)
        elif ext in ("png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp"):
            return _extraire_image(data)
        elif ext in ("txt", "csv", "rtf", "md"):
            return _extraire_txt(data)
        else:
            # Tentative texte brut en dernier recours
            try:
                return _extraire_txt(data)
            except Exception:
                return f"[Format .{ext} non supporté pour l'extraction de texte]"
    except Exception as e:
        st.warning(f"⚠️ Extraction partielle pour {uploaded_file.name} : {e}")
        return ""


def extraire_texte_multiple(uploaded_files) -> str:
    """
    Extrait et concatène le texte de plusieurs fichiers.
    Chaque fichier est séparé par un séparateur clair.
    Retourne le texte combiné tronqué à 12 000 caractères pour les LLMs.
    """
    if not uploaded_files:
        return ""

    parties = []
    for f in uploaded_files:
        texte = extraire_texte(f)
        if texte.strip():
            parties.append(
                f"━━━ DOCUMENT : {f.name} ({'%.1f' % (f.size/1024)} KB) ━━━\n{texte}"
            )

    combined = "\n\n".join(parties)
    # Tronquer intelligemment pour ne pas exploser le contexte LLM
    if len(combined) > 12000:
        combined = combined[:12000] + "\n\n[... Contenu tronqué à 12 000 caractères]"

    return combined


def feedback_fichiers(uploaded_files) -> None:
    """Affiche un résumé visuel des fichiers sélectionnés."""
    if not uploaded_files:
        st.caption("_Aucun document sélectionné_")
        return

    total_kb = sum(f.size for f in uploaded_files) / 1024
    size_label = f"{total_kb/1024:.1f} MB" if total_kb >= 1024 else f"{total_kb:.0f} KB"
    st.success(
        f"✅ **{len(uploaded_files)} fichier(s) prêt(s)** — "
        f"{', '.join(f.name for f in uploaded_files)} — {size_label} au total"
    )